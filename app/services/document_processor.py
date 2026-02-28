"""
Document Processing Service
Handles document parsing, text extraction, and intelligent chunking
"""
import os
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
import re

# Document parsers
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.logging_config import LoggerMixin, LogExecutionTime, get_logger
from app.core.exceptions import (
    DocumentParsingException,
    InvalidDocumentFormatException,
    DocumentTooLargeException
)
from app.core.retry_utils import retry_with_backoff, FallbackChain
from app.core.config import settings


logger = get_logger(__name__)


@dataclass
class DocumentChunk:
    """Represents a chunk of document text"""
    text: str
    chunk_index: int
    start_char: int
    end_char: int
    metadata: Dict[str, Any]
    token_count: Optional[int] = None


@dataclass
class ParsedDocument:
    """Represents a parsed document"""
    text: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any]
    total_chars: int
    total_chunks: int


class DocumentProcessor(LoggerMixin):
    """
    Service for processing documents
    Supports PDF, DOCX, and TXT files with intelligent chunking
    """
    
    def __init__(
        self,
        chunk_size: int = None,
        chunk_overlap: int = None,
        supported_formats: Optional[List[str]] = None,
        max_file_size_mb: Optional[int] = None
    ):
        """
        Initialize Document Processor
        
        Args:
            chunk_size: Size of text chunks in characters
            chunk_overlap: Overlap between chunks
            supported_formats: List of supported file extensions
            max_file_size_mb: Maximum file size in MB
        """
        self.chunk_size = chunk_size or settings.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP
        self.supported_formats = supported_formats or settings.get_supported_formats()
        self.max_file_size_mb = max_file_size_mb or settings.MAX_UPLOAD_SIZE_MB
        
        self.logger.info(
            f"Initialized Document Processor | "
            f"Chunk size: {self.chunk_size} | "
            f"Overlap: {self.chunk_overlap} | "
            f"Formats: {self.supported_formats}"
        )
    
    def process_file(
        self,
        file_path: str,
        document_metadata: Optional[Dict[str, Any]] = None
    ) -> ParsedDocument:
        """
        Process a document file
        
        Args:
            file_path: Path to document file
            document_metadata: Additional metadata for the document
            
        Returns:
            ParsedDocument with text and chunks
            
        Raises:
            InvalidDocumentFormatException: If file format not supported
            DocumentTooLargeException: If file exceeds size limit
            DocumentParsingException: If parsing fails
        """
        file_path = Path(file_path)
        
        # Validate file
        self._validate_file(file_path)
        
        self.logger.info(
            f"Processing document | "
            f"File: {file_path.name} | "
            f"Size: {file_path.stat().st_size / 1024:.1f} KB"
        )
        
        with LogExecutionTime(self.logger, f"Process document: {file_path.name}"):
            # Extract text based on file type
            file_ext = file_path.suffix.lower().lstrip('.')
            
            if file_ext == 'pdf':
                text = self._parse_pdf(file_path)
            elif file_ext == 'docx':
                text = self._parse_docx(file_path)
            elif file_ext == 'txt':
                text = self._parse_txt(file_path)
            elif file_ext == 'pptx':
                text = self._parse_pptx(file_path)
            else:
                raise InvalidDocumentFormatException(
                    filename=str(file_path),
                    file_type=file_ext,
                    supported_types=self.supported_formats
                )
            
            # Clean text
            text = self._clean_text(text)
            
            # Create chunks
            chunks = self._create_chunks(text, document_metadata or {})
            
            # Build metadata
            metadata = self._build_metadata(file_path, document_metadata)
            
            parsed_doc = ParsedDocument(
                text=text,
                chunks=chunks,
                metadata=metadata,
                total_chars=len(text),
                total_chunks=len(chunks)
            )
            
            self.logger.info(
                f"Processed document | "
                f"Chars: {len(text):,} | "
                f"Chunks: {len(chunks)}"
            )
            
            return parsed_doc
    
    def _validate_file(self, file_path: Path):
        """Validate file exists, format, and size"""
        # Check existence
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Check format
        file_ext = file_path.suffix.lower().lstrip('.')
        if file_ext not in self.supported_formats:
            raise InvalidDocumentFormatException(
                filename=str(file_path),
                file_type=file_ext,
                supported_types=self.supported_formats
            )
        
        # Check size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            raise DocumentTooLargeException(
                filename=str(file_path),
                size_mb=file_size_mb,
                max_size_mb=self.max_file_size_mb
            )
    
    def _parse_pdf(self, file_path: Path) -> str:
        """
        Parse PDF file with fallback strategies
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        def parse_with_pdfplumber():
            """Primary PDF parser using pdfplumber"""
            with pdfplumber.open(file_path) as pdf:
                text_parts = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                return "\n\n".join(text_parts)
        
        def parse_with_pypdf2():
            """Fallback PDF parser using PyPDF2"""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_parts = []
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                return "\n\n".join(text_parts)
        
        # Try parsers in order with fallback
        fallback_chain = FallbackChain(
            strategies=[parse_with_pdfplumber, parse_with_pypdf2],
            name="PDF Parsing"
        )
        
        try:
            return fallback_chain.execute()
        except Exception as e:
            raise DocumentParsingException(
                filename=str(file_path),
                reason="All PDF parsing strategies failed",
                original_exception=e
            )
    
    def _parse_docx(self, file_path: Path) -> str:
        """
        Parse DOCX file
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Extracted text
        """
        try:
            doc = DocxDocument(file_path)
            
            # Extract text from paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract text from tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        table_texts.append(row_text)
            
            # Combine paragraphs and tables
            all_text = "\n\n".join(paragraphs)
            if table_texts:
                all_text += "\n\n" + "\n".join(table_texts)
            
            return all_text
            
        except Exception as e:
            raise DocumentParsingException(
                filename=str(file_path),
                reason="Failed to parse DOCX file",
                original_exception=e
            )
            
    def _parse_pptx(self, file_path: Path) -> str:
        """
        Parse PPTX file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Extracted text from slides and shapes
        """
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    # Add a slide marker for better document context
                    text_parts.append(f"[Slide {slide_idx}]")
                    text_parts.append("\n".join(slide_text))
                    
            return "\n\n".join(text_parts)
            
        except Exception as e:
            raise DocumentParsingException(
                filename=str(file_path),
                reason="Failed to parse PPTX file",
                original_exception=e
            )

    def _parse_txt(self, file_path: Path) -> str:
        """
        Parse TXT file with encoding detection
        
        Args:
            file_path: Path to TXT file
            
        Returns:
            Extracted text
        """
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail
        raise DocumentParsingException(
            filename=str(file_path),
            reason="Could not decode file with any supported encoding",
            original_exception=None
        )
    
    def _clean_text(self, text: str) -> str:
        """
        Clean extracted text
        
        Args:
            text: Raw text
            
        Returns:
            Cleaned text
        """
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def _create_chunks(
        self,
        text: str,
        base_metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """
        Create overlapping chunks from text
        
        Args:
            text: Full document text
            base_metadata: Base metadata for all chunks
            
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        chunk_index = 0
        start_char = 0
        
        while start_char < len(text):
            # Calculate end position
            end_char = min(start_char + self.chunk_size, len(text))
            
            # Try to break at sentence boundary if not at document end
            if end_char < len(text):
                # Look for sentence-ending punctuation within last 20% of chunk
                search_start = int(end_char * 0.8)
                sentence_end = self._find_sentence_boundary(
                    text,
                    search_start,
                    end_char
                )
                
                if sentence_end > start_char:
                    end_char = sentence_end
            
            # Extract chunk text
            chunk_text = text[start_char:end_char].strip()
            
            if chunk_text:  # Only add non-empty chunks
                # Estimate token count (rough approximation: words * 1.3)
                word_count = len(chunk_text.split())
                token_count = int(word_count * 1.3)
                
                chunk = DocumentChunk(
                    text=chunk_text,
                    chunk_index=chunk_index,
                    start_char=start_char,
                    end_char=end_char,
                    metadata={
                        **base_metadata,
                        "chunk_index": chunk_index,
                        "word_count": word_count
                    },
                    token_count=token_count
                )
                
                chunks.append(chunk)
                chunk_index += 1
            
            # Move to next chunk with overlap
            start_char = end_char - self.chunk_overlap
            
            # Ensure we make progress
            if start_char <= chunks[-1].start_char if chunks else False:
                start_char = end_char
        
        return chunks
    
    def _find_sentence_boundary(
        self,
        text: str,
        start: int,
        end: int
    ) -> int:
        """
        Find sentence boundary within range
        
        Args:
            text: Full text
            start: Start position
            end: End position
            
        Returns:
            Position of sentence boundary, or end if not found
        """
        # Look for sentence-ending punctuation
        sentence_enders = ['. ', '! ', '? ', '.\n', '!\n', '?\n']
        
        best_pos = end
        for ender in sentence_enders:
            pos = text.rfind(ender, start, end)
            if pos != -1:
                # Found sentence ender, return position after it
                best_pos = pos + len(ender)
                break
        
        return best_pos
    
    def _build_metadata(
        self,
        file_path: Path,
        custom_metadata: Optional[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """Build comprehensive metadata for document"""
        metadata = {
            "filename": file_path.name,
            "file_path": str(file_path),
            "file_extension": file_path.suffix.lstrip('.'),
            "file_size_bytes": file_path.stat().st_size,
            "file_size_kb": file_path.stat().st_size / 1024,
        }
        
        # Add custom metadata
        if custom_metadata:
            metadata.update(custom_metadata)
        
        return metadata
    
    def process_text_direct(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ParsedDocument:
        """
        Process text directly (without file)
        
        Args:
            text: Text to process
            metadata: Metadata for the text
            
        Returns:
            ParsedDocument
        """
        text = self._clean_text(text)
        chunks = self._create_chunks(text, metadata or {})
        
        return ParsedDocument(
            text=text,
            chunks=chunks,
            metadata=metadata or {},
            total_chars=len(text),
            total_chunks=len(chunks)
        )


# Singleton instance
_document_processor_instance: Optional[DocumentProcessor] = None


def get_document_processor() -> DocumentProcessor:
    """
    Get or create document processor singleton
    
    Returns:
        Document processor instance
    """
    global _document_processor_instance
    
    if _document_processor_instance is None:
        _document_processor_instance = DocumentProcessor()
    
    return _document_processor_instance


if __name__ == "__main__":
    # Test document processor
    processor = DocumentProcessor(chunk_size=500, chunk_overlap=50)
    
    # Test with sample text
    sample_text = """
    Machine Learning is a subset of Artificial Intelligence. 
    It focuses on building systems that can learn from data.
    Deep Learning is a subset of Machine Learning that uses neural networks.
    Natural Language Processing helps computers understand human language.
    """ * 10  # Repeat to create longer text
    
    parsed = processor.process_text_direct(
        text=sample_text,
        metadata={"topic": "ML", "type": "notes"}
    )
    
    print(f"Processed text:")
    print(f"   Total chars: {parsed.total_chars:,}")
    print(f"   Total chunks: {parsed.total_chunks}")
    print(f"\nFirst chunk:")
    print(f"   {parsed.chunks[0].text[:200]}...")
    print(f"   Tokens: {parsed.chunks[0].token_count}")
